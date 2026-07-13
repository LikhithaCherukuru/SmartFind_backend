from pptx import Presentation

from app.ai.extractors.base_extractor import BaseExtractor


class PPTXExtractor(BaseExtractor):

    def extract_text(self, file_path: str):

        presentation = Presentation(file_path)

        text = ""

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

        return text